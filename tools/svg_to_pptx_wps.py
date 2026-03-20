#!/usr/bin/env python3
"""
WPS-friendly editable PPTX exporter.

Converts a constrained SVG subset into native PPT elements so each object can
be selected and edited in WPS Office without relying on "Convert to Shapes".
"""

from __future__ import annotations

import argparse
import re
import shutil
import tempfile
import zipfile
from pathlib import Path
from typing import Iterable, List, Optional, Tuple
from xml.etree import ElementTree as ET

from PIL import Image
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_VERTICAL_ANCHOR
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Emu, Pt

from svg_to_pptx import create_notes_slide_rels_xml, create_notes_slide_xml, find_notes_files, markdown_to_plain_text

SVG_NS = "http://www.w3.org/2000/svg"
XLINK_NS = "http://www.w3.org/1999/xlink"
EMU_PER_PIXEL = 914400 / 96
DEFAULT_MARGIN = 8


def px(value: float | str | None) -> int:
    if value is None:
        return 0
    if isinstance(value, (int, float)):
        return int(round(value * EMU_PER_PIXEL))
    text = str(value).strip()
    if not text:
        return 0
    text = re.sub(r"px$", "", text)
    return int(round(float(text) * EMU_PER_PIXEL))


def px_float(value: float | str | None) -> float:
    if value is None:
        return 0.0
    if isinstance(value, (int, float)):
        return float(value)
    text = str(value).strip()
    if not text:
        return 0.0
    text = re.sub(r"px$", "", text)
    return float(text)


def parse_color(value: Optional[str]) -> Optional[RGBColor]:
    if not value or value in ("none", "transparent"):
        return None
    value = value.strip()
    if value.startswith("#"):
        hex_value = value[1:]
        if len(hex_value) == 3:
            hex_value = "".join(ch * 2 for ch in hex_value)
        if len(hex_value) == 6:
            return RGBColor.from_string(hex_value.upper())
    return None


def parse_font_weight(value: Optional[str]) -> bool:
    if not value:
        return False
    if value.lower() == "bold":
        return True
    try:
        return int(value) >= 600
    except ValueError:
        return False


def estimate_text_width_px(text: str, font_size: float, letter_spacing: float = 0.0) -> float:
    width = 0.0
    for ch in text:
        code = ord(ch)
        if ch == " ":
            width += font_size * 0.33
        elif 0x4E00 <= code <= 0x9FFF:
            width += font_size * 1.00
        elif ch in "MW@#%&":
            width += font_size * 0.90
        elif ch in "ilI1|":
            width += font_size * 0.32
        elif ch in ".,:;!`'":
            width += font_size * 0.28
        elif ch in "/\\-()[]{}":
            width += font_size * 0.40
        else:
            width += font_size * 0.58
    if len(text) > 1:
        width += letter_spacing * (len(text) - 1)
    return width


def text_lines_from_element(el: ET.Element) -> List[Tuple[str, float]]:
    tspans = [child for child in list(el) if child.tag == f"{{{SVG_NS}}}tspan"]
    if not tspans:
        return [((el.text or "").strip(), 0.0)] if (el.text or "").strip() else []
    lines: List[Tuple[str, float]] = []
    for tspan in tspans:
        line = "".join(tspan.itertext()).strip()
        if not line:
            continue
        dy = px_float(tspan.get("dy") or 0)
        lines.append((line, dy))
    return lines


def resolve_image_path(svg_path: Path, href: str) -> Optional[Path]:
    if not href:
        return None
    if href.startswith("data:"):
        return None
    path = (svg_path.parent / href).resolve()
    return path if path.exists() else None


def fit_image_box(box_w: float, box_h: float, img_w: float, img_h: float) -> Tuple[float, float]:
    ratio = min(box_w / img_w, box_h / img_h)
    return img_w * ratio, img_h * ratio


def add_rect(slide, el: ET.Element):
    x = px(el.get("x"))
    y = px(el.get("y"))
    w = px(el.get("width"))
    h = px(el.get("height"))
    rx = px_float(el.get("rx") or 0)
    shape_type = MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE if rx > 0 else MSO_AUTO_SHAPE_TYPE.RECTANGLE
    shape = slide.shapes.add_shape(shape_type, x, y, w, h)

    fill_color = parse_color(el.get("fill"))
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color

    stroke_color = parse_color(el.get("stroke"))
    if stroke_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = stroke_color
        shape.line.width = max(Emu(1), px(el.get("stroke-width") or 1))


def add_line(slide, el: ET.Element):
    x1 = px(el.get("x1"))
    y1 = px(el.get("y1"))
    x2 = px(el.get("x2"))
    y2 = px(el.get("y2"))
    line = slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT, x1, y1, x2, y2)
    stroke_color = parse_color(el.get("stroke")) or RGBColor.from_string("000000")
    line.line.color.rgb = stroke_color
    line.line.width = max(Emu(1), px(el.get("stroke-width") or 1))


def add_circle(slide, el: ET.Element):
    cx = px_float(el.get("cx"))
    cy = px_float(el.get("cy"))
    r = px_float(el.get("r"))
    shape = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.OVAL,
        px(cx - r),
        px(cy - r),
        px(r * 2),
        px(r * 2),
    )
    fill_color = parse_color(el.get("fill"))
    if fill_color is None:
        shape.fill.background()
    else:
        shape.fill.solid()
        shape.fill.fore_color.rgb = fill_color
    stroke_color = parse_color(el.get("stroke"))
    if stroke_color is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = stroke_color
        shape.line.width = max(Emu(1), px(el.get("stroke-width") or 1))


def parse_points(points_text: str) -> List[Tuple[float, float]]:
    parts = re.split(r"[\s,]+", points_text.strip())
    nums = [float(p) for p in parts if p]
    return list(zip(nums[::2], nums[1::2]))


def add_polyline(slide, el: ET.Element):
    points = parse_points(el.get("points") or "")
    if len(points) < 2:
        return
    start_x, start_y = points[0]
    builder = slide.shapes.build_freeform(start_x=px(start_x), start_y=px(start_y), scale=1.0)
    builder.add_line_segments([(px(x), px(y)) for x, y in points[1:]], close=False)
    shape = builder.convert_to_shape()
    shape.fill.background()
    stroke_color = parse_color(el.get("stroke")) or RGBColor.from_string("000000")
    shape.line.color.rgb = stroke_color
    shape.line.width = max(Emu(1), px(el.get("stroke-width") or 1))


def add_image(slide, el: ET.Element, svg_path: Path):
    href = el.get("href") or el.get(f"{{{XLINK_NS}}}href")
    image_path = resolve_image_path(svg_path, href or "")
    if image_path is None:
        return

    x = px_float(el.get("x"))
    y = px_float(el.get("y"))
    w = px_float(el.get("width"))
    h = px_float(el.get("height"))
    preserve = el.get("preserveAspectRatio") or ""

    with Image.open(image_path) as img:
        img_w, img_h = img.size

    if "meet" in preserve:
        draw_w, draw_h = fit_image_box(w, h, img_w, img_h)
        left = x + (w - draw_w) / 2
        top = y + (h - draw_h) / 2
        slide.shapes.add_picture(str(image_path), px(left), px(top), px(draw_w), px(draw_h))
        return

    picture = slide.shapes.add_picture(str(image_path), px(x), px(y), px(w), px(h))
    if "slice" in preserve:
        img_ratio = img_w / img_h
        box_ratio = w / h if h else img_ratio
        if img_ratio > box_ratio:
            crop = (1 - (box_ratio / img_ratio)) / 2
            picture.crop_left = crop
            picture.crop_right = crop
        elif img_ratio < box_ratio:
            crop = (1 - (img_ratio / box_ratio)) / 2
            picture.crop_top = crop
            picture.crop_bottom = crop


def add_text(slide, el: ET.Element, slide_w_px: float):
    lines = text_lines_from_element(el)
    if not lines:
        return

    font_size = px_float(el.get("font-size") or 18)
    font_name = el.get("font-family") or "Microsoft YaHei"
    italic = (el.get("font-style") or "").lower() == "italic"
    bold = parse_font_weight(el.get("font-weight"))
    color = parse_color(el.get("fill")) or RGBColor.from_string("1F1A1C")
    text_anchor = (el.get("text-anchor") or "start").lower()
    letter_spacing = px_float(el.get("letter-spacing") or 0)
    x = px_float(el.get("x") or 0)
    y = px_float(el.get("y") or 0)

    width_px = max(estimate_text_width_px(line, font_size, letter_spacing) for line, _ in lines) + DEFAULT_MARGIN * 2
    line_height_px = font_size * 1.35
    height_px = line_height_px * len(lines) + DEFAULT_MARGIN * 2

    if text_anchor == "middle":
        left_px = x - width_px / 2
    elif text_anchor == "end":
        left_px = x - width_px
    else:
        left_px = x
    left_px = max(0, left_px)
    if left_px + width_px > slide_w_px:
        width_px = max(40, slide_w_px - left_px)

    top_px = max(0, y - font_size * 0.9)

    box = slide.shapes.add_textbox(px(left_px), px(top_px), px(width_px), px(height_px))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = False
    tf.margin_left = 0
    tf.margin_right = 0
    tf.margin_top = 0
    tf.margin_bottom = 0
    tf.vertical_anchor = MSO_VERTICAL_ANCHOR.TOP

    if text_anchor == "middle":
        align = PP_ALIGN.CENTER
    elif text_anchor == "end":
        align = PP_ALIGN.RIGHT
    else:
        align = PP_ALIGN.LEFT

    for idx, (line_text, _) in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = 0
        p.space_before = 0
        run = p.add_run()
        run.text = line_text
        run.font.name = font_name
        run.font.size = Pt(font_size)
        run.font.bold = bold
        run.font.italic = italic
        run.font.color.rgb = color
        force_east_asian_font(run, font_name)


def force_east_asian_font(run, font_name: str):
    """
    Explicitly set Latin/East Asian/Complex Script typefaces in OOXML.

    WPS often ignores `run.font.name` for Chinese text unless `a:ea` is set.
    """
    r_pr = run._r.get_or_add_rPr()
    for tag in ("a:latin", "a:ea", "a:cs"):
        child = r_pr.find(tag, {"a": "http://schemas.openxmlformats.org/drawingml/2006/main"})
        if child is None:
            child = OxmlElement(tag)
            r_pr.append(child)
        child.set("typeface", font_name)


def iter_svg_elements(root: ET.Element) -> Iterable[ET.Element]:
    for child in list(root):
        yield child


def viewbox_size(root: ET.Element) -> Tuple[float, float]:
    vb = root.get("viewBox")
    if vb:
        parts = re.split(r"[\s,]+", vb.strip())
        if len(parts) >= 4:
            return float(parts[2]), float(parts[3])
    width = px_float(root.get("width") or 1280)
    height = px_float(root.get("height") or 720)
    return width, height


def render_svg_to_slide(slide, svg_path: Path):
    root = ET.parse(svg_path).getroot()
    slide_w_px, _ = viewbox_size(root)
    for el in iter_svg_elements(root):
        tag = el.tag.split("}")[-1]
        if tag == "rect":
            add_rect(slide, el)
        elif tag == "line":
            add_line(slide, el)
        elif tag == "circle":
            add_circle(slide, el)
        elif tag == "polyline":
            add_polyline(slide, el)
        elif tag == "image":
            add_image(slide, el, svg_path)
        elif tag == "text":
            add_text(slide, el, slide_w_px)


def get_slide_dimensions_from_svg(svg_path: Path) -> Tuple[int, int, float, float]:
    root = ET.parse(svg_path).getroot()
    width_px, height_px = viewbox_size(root)
    return px(width_px), px(height_px), width_px, height_px


def embed_notes_in_pptx(pptx_path: Path, notes: dict[str, str], slide_count: int):
    if not notes:
        return
    temp_dir = Path(tempfile.mkdtemp())
    try:
        with zipfile.ZipFile(pptx_path, "r") as zf:
            zf.extractall(temp_dir)

        notes_slides_dir = temp_dir / "ppt" / "notesSlides"
        notes_slides_dir.mkdir(exist_ok=True)
        notes_rels_dir = notes_slides_dir / "_rels"
        notes_rels_dir.mkdir(exist_ok=True)

        for slide_num in range(1, slide_count + 1):
            slide_xml = temp_dir / "ppt" / "slides" / f"slide{slide_num}.xml"
            if not slide_xml.exists():
                continue
            stem = f"{slide_num:02d}"
            matched_key = ""
            for key in notes:
                if key.startswith(stem):
                    matched_key = key
                    break
            notes_text = markdown_to_plain_text(notes.get(matched_key, ""))

            notes_xml_path = notes_slides_dir / f"notesSlide{slide_num}.xml"
            notes_xml_path.write_text(create_notes_slide_xml(slide_num, notes_text), encoding="utf-8")
            notes_rels_path = notes_rels_dir / f"notesSlide{slide_num}.xml.rels"
            notes_rels_path.write_text(create_notes_slide_rels_xml(slide_num), encoding="utf-8")

            rels_path = temp_dir / "ppt" / "slides" / "_rels" / f"slide{slide_num}.xml.rels"
            content = rels_path.read_text(encoding="utf-8")
            if "notesSlide" not in content:
                rel = (
                    f'  <Relationship Id="rId10" '
                    f'Type="http://schemas.openxmlformats.org/officeDocument/2006/relationships/notesSlide" '
                    f'Target="../notesSlides/notesSlide{slide_num}.xml"/>'
                )
                content = content.replace("</Relationships>", rel + "\n</Relationships>")
                rels_path.write_text(content, encoding="utf-8")

        content_types_path = temp_dir / "[Content_Types].xml"
        content_types = content_types_path.read_text(encoding="utf-8")
        for slide_num in range(1, slide_count + 1):
            override = (
                f'  <Override PartName="/ppt/notesSlides/notesSlide{slide_num}.xml" '
                f'ContentType="application/vnd.openxmlformats-officedocument.presentationml.notesSlide+xml"/>'
            )
            if override not in content_types:
                content_types = content_types.replace("</Types>", override + "\n</Types>")
        content_types_path.write_text(content_types, encoding="utf-8")

        rebuilt = pptx_path.with_suffix(".tmp.pptx")
        with zipfile.ZipFile(rebuilt, "w", zipfile.ZIP_DEFLATED) as zf:
            for file_path in temp_dir.rglob("*"):
                if file_path.is_file():
                    zf.write(file_path, file_path.relative_to(temp_dir))
        rebuilt.replace(pptx_path)
    finally:
        shutil.rmtree(temp_dir, ignore_errors=True)


def find_svg_files(project_path: Path, source: str) -> List[Path]:
    source_dir = {
        "output": "svg_output",
        "final": "svg_final",
    }.get(source, source)
    svg_dir = project_path / source_dir
    return sorted(svg_dir.glob("*.svg"))


def default_output_path(project_path: Path) -> Path:
    return project_path / f"{project_path.name}_wps_editable.pptx"


def main():
    parser = argparse.ArgumentParser(description="Export SVG deck to WPS-editable PPTX")
    parser.add_argument("project_path", help="Project directory")
    parser.add_argument("-s", "--source", default="output", help="SVG source dir: output/final")
    parser.add_argument("-o", "--output", help="Output pptx path")
    parser.add_argument("--no-notes", action="store_true", help="Do not embed speaker notes")
    args = parser.parse_args()

    project_path = Path(args.project_path).resolve()
    svg_files = find_svg_files(project_path, args.source)
    if not svg_files:
        raise SystemExit("No SVG files found.")

    output_path = Path(args.output).resolve() if args.output else default_output_path(project_path)
    slide_w, slide_h, _, _ = get_slide_dimensions_from_svg(svg_files[0])

    prs = Presentation()
    prs.slide_width = slide_w
    prs.slide_height = slide_h
    blank = prs.slide_layouts[6]

    for svg_path in svg_files:
        slide = prs.slides.add_slide(blank)
        render_svg_to_slide(slide, svg_path)

    prs.save(str(output_path))

    if not args.no_notes:
        notes = find_notes_files(project_path, svg_files)
        embed_notes_in_pptx(output_path, notes, len(svg_files))

    print(f"[OK] WPS editable PPTX saved to: {output_path}")
    print(f"Slides: {len(svg_files)}")


if __name__ == "__main__":
    main()
